
from __future__ import unicode_literals

from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LABEL_POSITION,
    XL_LEGEND_POSITION,
    XL_TICK_MARK,
    XL_TICK_LABEL_POSITION
    )

from pptx.util import (
    Emu,
    Pt,
    Cm
    )

from pptx.enum.dml import (
    MSO_THEME_COLOR,
    MSO_COLOR_TYPE,
    MSO_FILL
    )

from pptx.enum.text import (
    PP_ALIGN,
    MSO_AUTO_SIZE,
    MSO_ANCHOR
    )

from os import path
from quantipy.core.builds.powerpoint.transformations import color_setter

thisdir = '%s/fonts/' % (path.split(__file__)[0])

"""  Simplified access to, and manipulation of python-pptx shapes/objects/properties.
"""

#------------------------------------------------------------------------------
#-------------------------------------------------------------------------------

fill_type_dct = {'background': MSO_FILL.BACKGROUND,
                 'gradient': MSO_FILL.GRADIENT,
                 'group': MSO_FILL.GROUP,
                 'patterned': MSO_FILL.PATTERNED,
                 'picture': MSO_FILL.PICTURE,
                 'solid': MSO_FILL.SOLID,
                 'textured': MSO_FILL.TEXTURED}

data_label_pos_dct = {'outside_end': XL_LABEL_POSITION.OUTSIDE_END,
                      'above': XL_LABEL_POSITION.ABOVE,
                      'below': XL_LABEL_POSITION.BELOW,
                      'best_fit': XL_LABEL_POSITION.BEST_FIT,
                      'center': XL_LABEL_POSITION.CENTER,
                      'inside_base': XL_LABEL_POSITION.INSIDE_BASE,
                      'inside_end': XL_LABEL_POSITION.INSIDE_END,
                      'left': XL_LABEL_POSITION.LEFT,
                      'mixed': XL_LABEL_POSITION.MIXED,
                      'outside_end': XL_LABEL_POSITION.OUTSIDE_END,
                      'right': XL_LABEL_POSITION.RIGHT}

legend_pos_dct = {'bottom': XL_LEGEND_POSITION.BOTTOM,
                  'corner': XL_LEGEND_POSITION.CORNER,
                  'custom': XL_LEGEND_POSITION.CUSTOM,
                  'left': XL_LEGEND_POSITION.LEFT,
                  'right': XL_LEGEND_POSITION.RIGHT,
                  'top': XL_LEGEND_POSITION.TOP}


tick_label_pos_dct = {'high': XL_TICK_LABEL_POSITION.HIGH,
                      'low': XL_TICK_LABEL_POSITION.LOW,
                      'next_to_axis': XL_TICK_LABEL_POSITION.NEXT_TO_AXIS,
                      'none': XL_TICK_LABEL_POSITION.NONE}


tick_mark_pos_dct = {'cross': XL_TICK_MARK.CROSS,
                     'inside': XL_TICK_MARK.INSIDE,
                     'none': XL_TICK_MARK.NONE,
                     'outside': XL_TICK_MARK.OUTSIDE}

vertical_alignment_pos_dct = {'top': MSO_ANCHOR.TOP,
                              'middle': MSO_ANCHOR.MIDDLE,
                              'bottom': MSO_ANCHOR.BOTTOM,
                              'mixed': MSO_ANCHOR.MIXED}


paragraph_alignment_pos_dct = {'center': PP_ALIGN.CENTER,
                               'distribute': PP_ALIGN.DISTRIBUTE,
                               'justify': PP_ALIGN.JUSTIFY,
                               'justify_low': PP_ALIGN.JUSTIFY_LOW,
                               'left': PP_ALIGN.LEFT,
                               'right': PP_ALIGN.RIGHT,
                               'thai_distribute': PP_ALIGN.THAI_DISTRIBUTE,
                               'mixed': PP_ALIGN.MIXED}


theme_color_index_dct = {'not_theme_color': MSO_THEME_COLOR.NOT_THEME_COLOR,
                         'accent_1': MSO_THEME_COLOR.ACCENT_1,
                         'accent_2': MSO_THEME_COLOR.ACCENT_2,
                         'accent_3': MSO_THEME_COLOR.ACCENT_3,
                         'accent_4': MSO_THEME_COLOR.ACCENT_4,
                         'accent_5': MSO_THEME_COLOR.ACCENT_5,
                         'accent_6': MSO_THEME_COLOR.ACCENT_6,
                         'background_1': MSO_THEME_COLOR.BACKGROUND_1,
                         'background_2': MSO_THEME_COLOR.BACKGROUND_2,
                         'dark_1': MSO_THEME_COLOR.DARK_1,
                         'dark_2': MSO_THEME_COLOR.DARK_2,
                         'followed_hyperlink': MSO_THEME_COLOR.FOLLOWED_HYPERLINK,
                         'hyperlink': MSO_THEME_COLOR.HYPERLINK,
                         'light_1': MSO_THEME_COLOR.LIGHT_1,
                         'light_2': MSO_THEME_COLOR.LIGHT_2,
                         'text_1': MSO_THEME_COLOR.TEXT_1,
                         'text_2': MSO_THEME_COLOR.TEXT_2,
                         'mixed': MSO_THEME_COLOR.MIXED}

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def percentage_of_num(percent, whole):
    """returns percent of a number. e.g. what is 5% of 20
    """

    return (percent * whole) / 100.0

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_cht_plot_height(chart_height, percent=10.52395879982087):
    """calculates a given chart's plot height by the charts height
    """

    amounttoremove = percentage_of_num(percent, chart_height)
    plot_height = chart_height - amounttoremove

    return int(plot_height)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def get_upper_cht_plot_gap(chart_height, percent=3.5078369905956115):
    """calculates the gap between the top of the plot area and top of chart areas
    """

    upper_cht_plot_gap = percentage_of_num(percent, chart_height)

    return int(upper_cht_plot_gap)

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_textbox(
              sld, text,
              left=0, top=0, width=300, height=100,
              font_name="Calibri",
              font_size=12,
              font_bold=True,
              font_italic=False,
              font_color=(89,89,89),
              font_color_brightness=0,
              font_color_theme=None,
              word_wrap=True,
              auto_size=None,
              fit_text=True,
              font_file=None,
              margin_left=0.25,
              margin_right=0.25,
              margin_top=0.13,
              margin_bottom=0.13,
              vertical_alignment='top',
              horizontal_alignment='left',
              textbox_fill_solid=False,
              textbox_color=(100,0,0),
              textbox_color_brightness=0,
              ):
    #-------------------------------------------------------------------------

    """Adds textbox and sets it's properties to a given slide.

    The only required arguments are the slide and the actual text
    content.

    Optional arguments:
    ________________________________________________________________________________________________________
    left                           |number of points the left edge of the textbox should
                                   |be from the left edge of the slide (default: 0)
    ________________________________________________________________________________________________________
    top                            |number of points the top edge of the textbox should be
                                   |from the top edge of the slide (default: 0)
    ________________________________________________________________________________________________________
    width                          |width of the textbox in points (default: 300)
    ________________________________________________________________________________________________________
    height                         |height of the textbox in points (default: 100)
    ________________________________________________________________________________________________________
    font_name                      |sets font type
    ________________________________________________________________________________________________________
    font_size                      |font size (default: 12)
    ________________________________________________________________________________________________________
    font_bold                      |turns bold on/off
    ________________________________________________________________________________________________________
    font_italic                    |turns italics on/off
    ________________________________________________________________________________________________________
    font_color_brightness          |range between -1.0 (darker) to 1.0 (lighter)
    ________________________________________________________________________________________________________
    word_wrap                      |turns word wrapping on/off in the new rectangle (default: True)
    ________________________________________________________________________________________________________
    auto_size                      |determines whether or not to automatically adjust the
                                   |height of the textbox to fit all the text
                                   |(default: 0 [do no auto-size, NONE, SHAPE_TO_FIT_TEXT, TEXT_TO_FIT_SHAPE])
    ________________________________________________________________________________________________________
    alignment                      |text alignment (default: 1 [left; see PowerPoint's
                                   |PpParagraphAlignment]
    ________________________________________________________________________________________________________
    vertical_alignment             |vertical text alignment
                                   |(default: 3 [middle; see Office's MsoVerticalAnchor])
    ________________________________________________________________________________________________________

    """

    textbox = sld.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))

    textframe = textbox.text_frame
    textframe.vertical_anchor = vertical_alignment_pos_dct[vertical_alignment]

    textframe.margin_left = Cm(margin_left)
    textframe.margin_bottom = Cm(margin_bottom)
    textframe.margin_right = Cm(margin_right)
    textframe.margin_top = Cm(margin_top)

    paragraph = textframe.paragraphs[0]
    paragraph.font.color.rgb = RGBColor(*font_color)
    if font_color_theme is not None:
        paragraph.font.color.theme_color = theme_color_index_dct[font_color_theme]
    paragraph.font.color.brightness = font_color_brightness
    paragraph.alignment = paragraph_alignment_pos_dct[horizontal_alignment]

    textframe.text = text
    if fit_text == True:
        if font_name == "Calibri":
            calibriz = path.join(thisdir, 'calibriz.ttf')
            textframe.fit_text(font_family=font_name, max_size=font_size, bold=font_bold,
                               italic=font_italic, font_file=calibriz)
        else:
            textframe.fit_text(font_family=font_name, max_size=font_size, bold=font_bold,
                               italic=font_italic)
    else:
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = font_bold
        paragraph.font.italic = font_italic
        textframe.word_wrap = word_wrap
        if auto_size is not None:
            textframe.auto_size = auto_size

    if textbox_fill_solid:
        txfill = textbox.fill
        txfill.solid()
        txfill.fore_color.rgb = RGBColor(*textbox_color)
        txfill.fore_color.brightness = textbox_color_brightness

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_pie_chart(
                sld, dataframe,
                left=284400, top=1475999, width=8582400, height=4140000,
                chart_style=2,

                #Legend properties
                has_legend=True,
                legend_position='bottom',
                legend_in_layout=False,
                legend_horz_offset=0,
                legend_font_name="Calibri",
                legend_font_size=10,
                legend_font_bold=False,
                legend_font_italic=False,
                legend_font_color=(89,89,89),
                legend_font_brightness=0,

                #Datalabel properties
                plot_has_data_labels=True,
                data_labels_position='outside_end',
                data_labels_num_format='0%',
                data_labels_num_format_is_linked=False,
                data_labels_font_name="Calibri",
                data_labels_font_size=9,
                data_labels_font_bold=False,
                data_labels_font_italic=False,
                data_labels_font_color=(0,0,0),
                ):
    #-------------------------------------------------------------------------
    """Adds single series pie chart to a given slide and set it's properties.

    The only required arguments are 'sld' and 'dataframe' content.

    Optional arguments:
    ________________________________________________________________________________________________________
    left                             |The distance, in EMU points, from left edge of Slide area.
                                     |to the left edge of the chart area. (default: 144000).
    ________________________________________________________________________________________________________
    top                              |The distance, in EMU points, from the top edge of the
                                     |Slide area to the top of the chart area. (default: 1584000).
    ________________________________________________________________________________________________________
    width                            |Width of the chart in EMU points (default: 8838000).
                                     |
    ________________________________________________________________________________________________________
    height                           |Height of the chart in EMU points (default: 4320000).
                                     |
    ________________________________________________________________________________________________________
    chart_style                      |Sets the chart style for the chart.
                                     |You can use a number from 1 to 48 to set the chart style (default: 2).
    ________________________________________________________________________________________________________
    has_legend                       |Boolean. True if the chart has a legend.
                                     |
    ________________________________________________________________________________________________________
    legend_position                  |Sets the position of the legend on the chart.
                                     |[bottom, corner, custom, left, right, top] (default: bottom)
    ________________________________________________________________________________________________________
    legend_in_layout                 |Boolean. True if a legend will occupy the chart layout space when
                                     |a chart layout is being determined.
    ________________________________________________________________________________________________________
    legend_horz_offset               |Adjustment of the x position of the legend from its default. Expressed
                                     |as a float between -1.0 and 1.0 representing a fraction of the chart width.
    ________________________________________________________________________________________________________
    legend_font_name                 |Set the typeface name for this Font instance, causing the text it
                                     |controls to appear in the named font, if a matching font is found.
    ________________________________________________________________________________________________________
    legend_font_size                 |Set length value or None, indicating the font height
                                     |in EMU.
    ________________________________________________________________________________________________________
    legend_font_bold                 |Boolean. Set boolean bold value of Font
                                     |
    ________________________________________________________________________________________________________
    legend_font_italic               |Boolean. Set boolean italic value of Font
                                     |
    ________________________________________________________________________________________________________
    legend_font_color                |Set color for this lengend's font in RGBColor(r, g, b)
                                     |
    ________________________________________________________________________________________________________
    legend_font_brightness           |Set float value between -1.0 and 1.0 indicating the brightness adjustment
                                     |for this color, e.g. -0.25 is 25 percent darker and 0.4 is 40 percent lighter.
                                     |0 means no brightness adjustment.
    ________________________________________________________________________________________________________
    plot_has_data_labels             |Boolean. Set if the series has data labels.
                                     |Assigning True causes data labels to be added to the plot
    ________________________________________________________________________________________________________
    data_labels_position             |value specifying the position of the data labels
                                     |with respect to their data point, or None if no position is specified
                                     |[outside_end, above, below, best_fit, center, inside_base, inside_end,
                                     |left, mixed, outside_end, right] (default: outside_end).
    ________________________________________________________________________________________________________
    data_labels_num_format_is_linked |Set string specifying the format for the numbers on this set of data labels
                                     |Returns 'General' if no number format has been set. Note that this format string
                                     |has no effect on rendered data labels when number_format_is_linked() is True
    ________________________________________________________________________________________________________
    data_labels_font_name            |Set the typeface name for this Font instance, causing the text it
                                     |controls to appear in the named font, if a matching font is found.
    ________________________________________________________________________________________________________
    data_labels_font_size            |Set length value or None, indicating the font height
                                     |in EMU.
    ________________________________________________________________________________________________________
    data_labels_font_bold            |Boolean. Set boolean bold value of Font
                                     |
    ________________________________________________________________________________________________________
    data_labels_font_italic          |Boolean. Set boolean italic value of Font
                                     |
    ________________________________________________________________________________________________________
    data_labels_font_color           |Set color for this data lebel's font in RGBColor(r, g, b)
                                     |
    ________________________________________________________________________________________________________

    """

    # Adding chart data
    chart_data = ChartData()
    chart_data.categories = dataframe.index

    for i, col in enumerate(dataframe.columns):
        chart_data.add_series(col, (dataframe.ix[:, i].values), '0.00%')

    # Adding chart
    x, y, cx, cy = Emu(left), Emu(top), Emu(width), Emu(height)
    graphic_frame = sld.shapes.add_chart(
        XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    #---------------- adjust chart properties -------------------------------

    chart.chart_style = chart_style

    # set legend properties
    chart.has_legend = has_legend
    if has_legend:
        legend = chart.legend
        legend.font.name = legend_font_name
        legend.font.size = Pt(legend_font_size)
        legend.font.bold = legend_font_bold
        legend.font.italic = legend_font_italic
        legend.font.color.rgb  = RGBColor(*legend_font_color)
        legend.font.color.brightness = legend_font_brightness
        legend.position = legend_pos_dct[legend_position]
        legend.include_in_layout = legend_in_layout
        legend.horz_offset = legend_horz_offset

    plot = chart.plots[0]
    plot.has_data_labels = plot_has_data_labels

    # set datalabel properties
    if plot_has_data_labels:
        data_labels = plot.data_labels
        data_labels.position = data_label_pos_dct[data_labels_position]
        data_labels.font.name = data_labels_font_name
        data_labels.font.size = Pt(data_labels_font_size)
        data_labels.font.bold = data_labels_font_bold
        data_labels.font.italic = data_labels_font_italic
        data_labels.font.color.rgb = RGBColor(*data_labels_font_color)
        if data_labels_num_format is not None:
            data_labels.number_format = data_labels_num_format
        data_labels.number_format_is_linked = data_labels_num_format_is_linked

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_bar_chart(
                sld, dataframe,
                left=284400, top=1475999, width=8582400, height=4140000,
                chart_style=2,

                #Legend properties
                has_legend=False,
                legend_position='right',
                legend_in_layout=False,
                legend_horz_offset = 0.1583,
                legend_font_name="Calibri",
                legend_font_size=10,
                legend_font_bold=False,
                legend_font_italic=False,
                legend_font_color=(89,89,89),
                legend_font_brightness=0,

                #Category axis properties
                caxis_visible=True,
                caxis_tick_label_position='none',
                caxis_tick_labels_offset=730,
                caxis_has_major_gridlines=False,
                caxis_has_minor_gridlines=False,
                caxis_major_tick_mark='outside',
                caxis_minor_tick_mark='none',
                caxis_tick_labels_font_name="Calibri",
                caxis_tick_labels_font_size=10,
                caxis_tick_labels_font_bold=False,
                caxis_tick_labels_font_italic=False,
                caxis_tick_labels_font_color=(89,89,89),

                #Value axis properties
                vaxis_visible=True,
                vaxis_tick_label_position='low',
                vaxis_has_major_gridlines=True,
                vaxis_has_minor_gridlines=False,
                vaxis_major_tick_mark='outside',
                vaxis_minor_tick_mark='none',
                vaxis_max_scale=1,
                vaxis_min_scale=0,
                vaxis_major_unit=0.1,
                vaxis_minor_unit=None,
                vaxis_tick_labels_num_format='0%',
                vaxis_tick_labels_num_format_is_linked=False,
                vaxis_tick_labels_font_name="Calibri",
                vaxis_tick_labels_font_bold=True,
                vaxis_tick_labels_font_size=10,
                vaxis_tick_labels_font_italic=False,
                vaxis_tick_labels_font_color=(89,89,89),

                #Datalabel properties
                plot_has_data_labels=True,
                data_labels_position='outside_end',
                data_labels_num_format='0%',
                data_labels_num_format_is_linked=False,
                data_labels_font_name="Calibri",
                data_labels_font_size=9,
                data_labels_font_bold=False,
                data_labels_font_italic=False,
                data_labels_font_color=(0,0,0),

                #Plot properties
                plot_vary_by_cat=False,
                series_color_order='reverse',
                invert_series_color_if_negative=False,
                plot_gap_width=150,
                plot_overlap=-10
                ):
    #-------------------------------------------------------------------------
    """Adds single or multi series bar chart to a given slide and set it's properties.

    The only required arguments are slide and dataframe content.

    Optional arguments:
    ________________________________________________________________________________________________________
    left                           |number of points the left edge of the chart should
                                   |be from the left edge of the slide (default: 3500000)
    ________________________________________________________________________________________________________
    top                            |number of points the top edge of the chart should be
                                   |from the top edge of the slide (default: 1450000)
    ________________________________________________________________________________________________________
    width                          |width of the chart in points (default: 5000000)
    ________________________________________________________________________________________________________
    height                         |height of the chart in points (default: 4466000)
    ________________________________________________________________________________________________________
    chart_style                    |available options on the Chart Styles group on
                                   |the Design tab on the Ribbon (default: 2)
    ________________________________________________________________________________________________________
    has_legend                     |turns on/off the legend for chart (default: True)
    ________________________________________________________________________________________________________
    legend_in_layout               |determines if a legend will occupy the chart layout
                                   |space when a chart layout is being determined (default: False)
    ________________________________________________________________________________________________________
    legend_position                |determines the position of a legend
                                   |(default: XL_LEGEND_POSITION.RIGHT [TOP, BOTTOM,LEFT,RIGHT,TOP_RIGHT])
    ________________________________________________________________________________________________________
    legend_horz_offset             |sets the horizontal distance between the plot area and end of chart area
                                   |(default: 0.1583)
    ________________________________________________________________________________________________________
    caxis_visible                  |determines if the category labels are visible or not (default: False)
    ________________________________________________________________________________________________________
    caxis_has_major_gridlines      |turns major gridlines on/off for category axis
    ________________________________________________________________________________________________________
    caxis_has_minor_gridlines      |turns minor gridlines on/off for category axis
    ________________________________________________________________________________________________________
    caxis_major_tick_mark          |draws marks at the same point along an axis as a data point
                                   |(default: XL_TICK_MARK.NONE [OUTSIDE, INSIDE, CROSS, NONE])
    ________________________________________________________________________________________________________
    caxis_minor_tick_mark          |draws marks in between major tick marks, where labels usually occur
                                   |(default: XL_TICK_MARK.NONE [OUTSIDE, INSIDE, CROSS, NONE])
    ________________________________________________________________________________________________________
    caxis_tick_labels_offset       |label spacing of the category axis (default: 730)
    ________________________________________________________________________________________________________
    vaxis_visible                  |determines if the category labels are visible or not (default: True)
    ________________________________________________________________________________________________________
    vaxis_has_major_gridlines      |turns major gridlines on/off for value axis
    ________________________________________________________________________________________________________
    vaxis_has_minor_gridlines      |turns minor gridlines on/off for value axis
    ________________________________________________________________________________________________________
    vaxis_max_scale                |sets maximum values for the value axis (default: 100.0)
    ________________________________________________________________________________________________________
    vaxis_min_scale                |sets minimum values for the value axis (default: 0.0)
    ________________________________________________________________________________________________________
    vaxis_major_unit               |sets maximum values for the value axis (default: 10.0)
    ________________________________________________________________________________________________________
    vaxis_minor_unit               |sets minimum values for the value axis (default: None)
    ________________________________________________________________________________________________________
    vaxis_minor_tick_mark          |draws marks at the same point along an axis as a data point
                                   |(default: XL_TICK_MARK.NONE [OUTSIDE, INSIDE, CROSS, NONE])
    ________________________________________________________________________________________________________
    vaxis_major_tick_mark          |draws the minor tick marks for the value axis in Chart to be inside
                                   |the axis (default: XL_TICK_MARK.OUTSIDE, [OUTSIDE, INSIDE, CROSS, NONE])
    ________________________________________________________________________________________________________
    vaxis_num_format               |sets the number format for the tick-mark labels on the value axis in
                                   |the chart. (default: '0"%"')
    ________________________________________________________________________________________________________
    vaxis_font_bold                |sets font to bold on the value axis in chart (default: True)
    ________________________________________________________________________________________________________
    plot_vary_by_cat               |not sure this property does exactly, to find out.
    _______________________________________________________________________________________________________
    plot_has_data_labels           |enables data labels for series in chart (default: True)
    ________________________________________________________________________________________________________
    data_labels_position           |determines position of data labels
                                   |(default:XL_LABEL_POSITION.OUTSIDE_END [CENTER, INSIDE_END, INSIDE_BASE])
    ________________________________________________________________________________________________________
    data_label_num_format          |sets the number format for the data labels on series plot in
                                   |the chart. (default: None ['0"%"'])
    ________________________________________________________________________________________________________
    invert_if_negative             |inverts bar colour for negative series (default: False [same colour as rest])
    ________________________________________________________________________________________________________
    plot_gap_width                 |determine space between bars (default: 100 [range between 0 to 500)
    ________________________________________________________________________________________________________
    plot_overlap                   |You can set this property to a value from -100 through 100. If this
                                   |property is set to -100, bars are positioned so that there is one bar width
                                   |between them. If the overlap is 0 (zero), there is no space between bars
                                   |(one bar starts immediately after the preceding bar). If the overlap
                                   |is 100, bars are positioned on top of each other.
                                   |(default: -10)
    ________________________________________________________________________________________________________
    """

    #if category labels are split from the chart shape then determine the width of the textboxes and the chart shape.
    #textboxes in this case will take up 40% of the overall width of the chart shape. From this we can calculate the
    #width of the chart shape

    if (caxis_visible == False) or (caxis_visible == True and str(tick_label_pos_dct[caxis_tick_label_position]) == "NONE (-4142)"):
        catwidth = percentage_of_num(40, width)
        width = width - catwidth
        left = left + catwidth

    # orientation of chart type requires that we reverse the row and column order.
    dataframe = dataframe[::-1]
    dataframe = dataframe[dataframe.columns[::-1]]

    # add chart data
    chart_data = ChartData()
    chart_data.categories = dataframe.index

    for i, col in enumerate(dataframe.columns):
        chart_data.add_series(col, (dataframe.ix[:, i].values), '0.00%')

    # add chart
    x, y, cx, cy = Emu(left), Emu(top), Emu(width), Emu(height)
    graphic_frame = sld.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    # ---------------- adjust chart properties ----------------

    # chart style
    chart.chart_style = chart_style

    # set legend properties
    chart.has_legend = has_legend
    if has_legend:
        legend = chart.legend
        legend.include_in_layout = legend_in_layout
        legend.position = legend_pos_dct[legend_position]
        legend.horz_offset = legend_horz_offset
        legend.font.name = legend_font_name
        legend.font.size = Pt(legend_font_size)
        legend.font.bold = legend_font_bold
        legend.font.italic = legend_font_italic
        legend.font.color.rgb  = RGBColor(*legend_font_color)
        legend.font.color.brightness = legend_font_brightness

    # set category axis (vertical) properties
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = caxis_has_major_gridlines
    category_axis.has_minor_gridlines = caxis_has_minor_gridlines
    category_axis.major_tick_mark = tick_mark_pos_dct[caxis_major_tick_mark]
    category_axis.minor_tick_mark = tick_mark_pos_dct[caxis_minor_tick_mark]
    category_axis.tick_label_position = tick_label_pos_dct[caxis_tick_label_position]

    category_axis.visible = caxis_visible
    if caxis_visible:
        caxis_tick_labels = category_axis.tick_labels
        caxis_tick_labels.offset = caxis_tick_labels_offset
        caxis_tick_labels.font.name = caxis_tick_labels_font_name
        caxis_tick_labels.font.size = Pt(caxis_tick_labels_font_size)
        caxis_tick_labels.font.bold = caxis_tick_labels_font_bold
        caxis_tick_labels.font.italic = caxis_tick_labels_font_italic
        caxis_tick_labels.font.color.rgb = RGBColor(*caxis_tick_labels_font_color)

    # set value axis (horizontal) properties
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = vaxis_has_major_gridlines
    value_axis.has_minor_gridlines = vaxis_has_minor_gridlines
    value_axis.maximum_scale = vaxis_max_scale
    value_axis.minimum_scale = vaxis_min_scale
    value_axis.major_unit = vaxis_major_unit
    value_axis.minor_unit = vaxis_minor_unit
    value_axis.major_tick_mark = tick_mark_pos_dct[vaxis_major_tick_mark]
    value_axis.minor_tick_mark = tick_mark_pos_dct[vaxis_minor_tick_mark]
    value_axis.tick_label_position = tick_label_pos_dct[vaxis_tick_label_position]

    value_axis.visible = vaxis_visible
    if vaxis_visible:
        vaxis_tick_labels = value_axis.tick_labels
        vaxis_tick_labels.font.bold = vaxis_tick_labels_font_bold
        vaxis_tick_labels.font.size = Pt(vaxis_tick_labels_font_size)
        vaxis_tick_labels.font.italic = vaxis_tick_labels_font_italic
        vaxis_tick_labels.font.name = vaxis_tick_labels_font_name
        vaxis_tick_labels.font.color.rgb = RGBColor(*vaxis_tick_labels_font_color)
        if vaxis_tick_labels_num_format is not None:
            vaxis_tick_labels.number_format = vaxis_tick_labels_num_format
        vaxis_tick_labels.number_format_is_linked = vaxis_tick_labels_num_format_is_linked

    # set plot area properties
    plot = chart.plots[0]
    plot.vary_by_categories = plot_vary_by_cat
    plot.gap_width = plot_gap_width
    plot.overlap = plot_overlap

    plot.has_data_labels = plot_has_data_labels
    if plot_has_data_labels:
        data_labels = plot.data_labels
        data_labels.position = data_label_pos_dct[data_labels_position]
        data_labels.font.size = Pt(data_labels_font_size)
        data_labels.font.bold = data_labels_font_bold
        data_labels.font.italic = data_labels_font_italic
        data_labels.font.name = data_labels_font_name
        data_labels.font.color.rgb = RGBColor(*data_labels_font_color)
        if data_labels_num_format is not None:
            data_labels.number_format = data_labels_num_format
        data_labels.number_format_is_linked = data_labels_num_format_is_linked

    if series_color_order and len(dataframe.columns) > 1:
        ser_colors_list = color_setter(len(dataframe.columns), series_color_order)

    for i, ser in enumerate(dataframe.columns):
        ser = plot.series[i]
        ser.invert_if_negative = invert_series_color_if_negative

        if series_color_order and len(dataframe.columns) > 1:
            fill = ser.fill
            fill.solid()
            color_code = ser_colors_list[i]
            fill.fore_color.rgb = RGBColor(*color_code)

    # generate overlay axis labels
    if (caxis_visible == False) or (caxis_visible == True and str(tick_label_pos_dct[caxis_tick_label_position]) == "NONE (-4142)"):
        cht_plot_height = get_cht_plot_height(height)
        heightPerLabel = cht_plot_height/len(dataframe.index)
        rightofchart = left + width
        txtbx_width = width / 5
        firstposition = top + get_upper_cht_plot_gap(height)

        cat_labels = dataframe.T.columns

        for i, label in enumerate(cat_labels):

            top = 0
            pointRelPos = len(cat_labels) - (i + 1)
            top = firstposition + pointRelPos * heightPerLabel

            add_textbox(sld,
                        left=142875, top=top, width=rightofchart - width, height=heightPerLabel,
                        text=label,
                        font_name=caxis_tick_labels_font_name,
                        font_size=caxis_tick_labels_font_size,
                        fit_text=False,
                        word_wrap=True,
                        font_bold=False,
                        font_color=caxis_tick_labels_font_color,
                        horizontal_alignment='right',
                        vertical_alignment='middle')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_column_chart(
                sld, dataframe,
                left=284400, top=1475999, width=8582400, height=4140000,
                chart_style=2,

                #Legend properties
                has_legend=True,
                legend_position='bottom',
                legend_in_layout=False,
                legend_horz_offset=0,
                legend_font_name="Calibri",
                legend_font_size=10,
                legend_font_bold=False,
                legend_font_italic=False,
                legend_font_color=(89,89,89),
                legend_font_brightness=0,

                #Category axis properties
                caxis_visible=True,
                caxis_tick_label_position='next_to_axis',
                caxis_tick_labels_offset=100,
                caxis_has_major_gridlines=False,
                caxis_has_minor_gridlines=False,
                caxis_major_tick_mark='outside',
                caxis_minor_tick_mark='none',
                caxis_tick_labels_font_name="Calibri",
                caxis_tick_labels_font_size=10,
                caxis_tick_labels_font_bold=False,
                caxis_tick_labels_font_italic=False,
                caxis_tick_labels_font_color=(89,89,89),

                #Value axis properties
                vaxis_visible=True,
                vaxis_tick_label_position='low',
                vaxis_has_major_gridlines=True,
                vaxis_has_minor_gridlines=False,
                vaxis_major_tick_mark='outside',
                vaxis_minor_tick_mark='none',
                vaxis_max_scale=1.0,
                vaxis_min_scale=0.0,
                vaxis_major_unit=0.1,
                vaxis_minor_unit=None,
                vaxis_tick_labels_num_format='0%',
                vaxis_tick_labels_num_format_is_linked=False,
                vaxis_tick_labels_font_name="Calibri",
                vaxis_tick_labels_font_size=10,
                vaxis_tick_labels_font_bold=True,
                vaxis_tick_labels_font_italic=False,
                vaxis_tick_labels_font_color=(89,89,89),

                #Datalabel properties
                plot_has_data_labels=True,
                data_labels_position='outside_end',
                data_labels_num_format='0%',
                data_labels_num_format_is_linked=False,
                data_labels_font_name="Calibri",
                data_labels_font_size=9,
                data_labels_font_bold=False,
                data_labels_font_italic=False,
                data_labels_font_color=(0,0,0),

                #Plot properties
                plot_vary_by_cat=False,
                invert_series_color_if_negative=False,
                plot_gap_width=150,
                plot_overlap=-10,
                ):
    #-------------------------------------------------------------------------

    # add chart data
    chart_data = ChartData()
    chart_data.categories = dataframe.index

    for col in dataframe.columns:
        chart_data.add_series(col, (dataframe[col].values), '0.00%')

    # add chart
    x, y, cx, cy = Emu(left), Emu(top), Emu(width), Emu(height)
    graphic_frame = sld.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    # ---------------- adjust chart properties ----------------

    # chart style
    chart.chart_style = chart_style

    # set legend properties
    chart.has_legend = has_legend
    if has_legend:
        legend = chart.legend
        legend.include_in_layout = legend_in_layout
        legend.position = legend_pos_dct[legend_position]
        legend.horz_offset = legend_horz_offset
        legend.font.name = legend_font_name
        legend.font.size = Pt(legend_font_size)
        legend.font.bold = legend_font_bold
        legend.font.italic = legend_font_italic
        legend.font.color.rgb  = RGBColor(*legend_font_color)
        legend.font.color.brightness = legend_font_brightness

    # set category axis (horizontal) properties
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = caxis_has_major_gridlines
    category_axis.has_minor_gridlines = caxis_has_minor_gridlines
    category_axis.major_tick_mark = tick_mark_pos_dct[caxis_major_tick_mark]
    category_axis.minor_tick_mark = tick_mark_pos_dct[caxis_minor_tick_mark]
    category_axis.tick_label_position = tick_label_pos_dct[caxis_tick_label_position]

    category_axis.visible = caxis_visible
    if caxis_visible:
        caxis_tick_labels = category_axis.tick_labels
        caxis_tick_labels.offset = caxis_tick_labels_offset
        caxis_tick_labels.font.name = caxis_tick_labels_font_name
        caxis_tick_labels.font.size = Pt(caxis_tick_labels_font_size)
        caxis_tick_labels.font.bold = caxis_tick_labels_font_bold
        caxis_tick_labels.font.italic = caxis_tick_labels_font_italic
        caxis_tick_labels.font.color.rgb = RGBColor(*caxis_tick_labels_font_color)

    # set value axis (vertical) properties
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = vaxis_has_major_gridlines
    value_axis.has_minor_gridlines = vaxis_has_minor_gridlines
    value_axis.minor_tick_mark = tick_mark_pos_dct[vaxis_minor_tick_mark]
    value_axis.major_tick_mark = tick_mark_pos_dct[vaxis_major_tick_mark]
    value_axis.maximum_scale = vaxis_max_scale
    value_axis.minimum_scale = vaxis_min_scale
    value_axis.major_unit = vaxis_major_unit
    value_axis.minor_unit = vaxis_minor_unit
    value_axis.tick_label_position = tick_label_pos_dct[vaxis_tick_label_position]

    value_axis.visible = vaxis_visible
    if vaxis_visible:
        vaxis_tick_labels = value_axis.tick_labels
        vaxis_tick_labels.font.bold = vaxis_tick_labels_font_bold
        vaxis_tick_labels.font.size = Pt(vaxis_tick_labels_font_size)
        vaxis_tick_labels.font.italic = vaxis_tick_labels_font_italic
        vaxis_tick_labels.font.name = vaxis_tick_labels_font_name
        vaxis_tick_labels.font.color.rgb = RGBColor(*vaxis_tick_labels_font_color)
        if vaxis_tick_labels_num_format is not None:
            vaxis_tick_labels.number_format = vaxis_tick_labels_num_format
        vaxis_tick_labels.number_format_is_linked = vaxis_tick_labels_num_format_is_linked

    # set plot area properties
    plot = chart.plots[0]
    plot.vary_by_categories = plot_vary_by_cat
    plot.gap_width = plot_gap_width
    plot.overlap = plot_overlap

    plot.has_data_labels = plot_has_data_labels
    if plot_has_data_labels:
        data_labels = plot.data_labels
        data_labels.position = data_label_pos_dct[data_labels_position]
        data_labels.font.size = Pt(data_labels_font_size)
        data_labels.font.bold = data_labels_font_bold
        data_labels.font.italic = data_labels_font_italic
        data_labels.font.name = data_labels_font_name
        data_labels.font.color.rgb = RGBColor(*data_labels_font_color)
        if data_labels_num_format  is not None:
            data_labels.number_format = data_labels_num_format
        data_labels.number_format_is_linked = data_labels_num_format_is_linked

    for i, ser in enumerate(dataframe.columns):
        ser = plot.series[i]
        ser.invert_if_negative = invert_series_color_if_negative

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_line_chart(
                sld, dataframe,
                left=284400, top=1475999, width=8582400, height=4140000,
                chart_style=2,

                #Legend properties
                has_legend=True,
                legend_position='bottom',
                legend_in_layout=False,
                legend_horz_offset=0,
                legend_font_name="Calibri",
                legend_font_size=10,
                legend_font_bold=False,
                legend_font_italic=False,
                legend_font_color=(89,89,89),
                legend_font_brightness=0,

                #Category axis properties
                caxis_visible=True,
                caxis_tick_label_position='low',
                caxis_tick_labels_offset=100,
                caxis_has_major_gridlines=False,
                caxis_has_minor_gridlines=False,
                caxis_major_tick_mark='outside',
                caxis_minor_tick_mark='none',
                caxis_tick_labels_font_name="Calibri",
                caxis_tick_labels_font_size=10,
                caxis_tick_labels_font_italic=False,
                caxis_tick_labels_font_bold=False,
                caxis_tick_labels_font_color=(89,89,89),

                #Value axis properties
                vaxis_visible=True,
                vaxis_tick_label_position='low',
                vaxis_has_major_gridlines=True,
                vaxis_has_minor_gridlines=False,
                vaxis_major_tick_mark='outside',
                vaxis_minor_tick_mark='none',
                vaxis_max_scale=1.0,
                vaxis_min_scale=0.0,
                vaxis_major_unit=0.1,
                vaxis_minor_unit=None,
                vaxis_tick_labels_num_format='0%',
                vaxis_tick_labels_num_format_is_linked=False,
                vaxis_tick_labels_font_name="Calibri",
                vaxis_tick_labels_font_bold=True,
                vaxis_tick_labels_font_size=10,
                vaxis_tick_labels_font_italic=False,
                vaxis_tick_labels_font_color=(89,89,89),

                #Data label properties
                plot_has_data_labels=True,
                data_labels_position='above',
                data_labels_num_format='0%',
                data_labels_num_format_is_linked=False,
                data_labels_font_name="Calibri",
                data_labels_font_size=9,
                data_labels_font_bold=False,
                data_labels_font_italic=False,
                data_labels_font_color=(0,0,0),

                #Plot properties
                plot_vary_by_cat=False,
                invert_series_color_if_negative=False,
                plot_gap_width=150,
                plot_overlap=-10,
                smooth_line=False,
                ):
    #-------------------------------------------------------------------------

    # Adding chart data
    chart_data = ChartData()
    chart_data.categories = dataframe.index

    for i, col in enumerate(dataframe.columns):
        chart_data.add_series(col, (dataframe.ix[:, i].values), '0.00%')

    # Adding chart
    x, y, cx, cy = Emu(left), Emu(top), Emu(width), Emu(height)
    graphic_frame = sld.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    # ---------------- adjust chart properties ----------------

    #chart style
    chart.chart_style = chart_style

    #set legend properties
    chart.has_legend = has_legend
    if has_legend:
        legend = chart.legend
        legend.include_in_layout = legend_in_layout
        legend.position = legend_pos_dct[legend_position]
        legend.horz_offset = legend_horz_offset
        legend.font.name = legend_font_name
        legend.font.size = Pt(legend_font_size)
        legend.font.bold = legend_font_bold
        legend.font.italic = legend_font_italic
        legend.font.color.rgb  = RGBColor(*legend_font_color)
        legend.font.color.brightness = legend_font_brightness

    #set category axis (horizontal) properties
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = caxis_has_major_gridlines
    category_axis.has_minor_gridlines = caxis_has_minor_gridlines
    category_axis.major_tick_mark = tick_mark_pos_dct[caxis_major_tick_mark]
    category_axis.minor_tick_mark = tick_mark_pos_dct[caxis_minor_tick_mark]
    category_axis.tick_label_position = tick_label_pos_dct[caxis_tick_label_position]

    category_axis.visible = caxis_visible
    if caxis_visible:
        caxis_tick_labels = category_axis.tick_labels
        caxis_tick_labels.offset = caxis_tick_labels_offset
        caxis_tick_labels.font.name = caxis_tick_labels_font_name
        caxis_tick_labels.font.size = Pt(caxis_tick_labels_font_size)
        caxis_tick_labels.font.bold = caxis_tick_labels_font_bold
        caxis_tick_labels.font.italic = caxis_tick_labels_font_italic
        caxis_tick_labels.font.color.rgb = RGBColor(*caxis_tick_labels_font_color)

    # set value axis (vertical) properties
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = vaxis_has_major_gridlines
    value_axis.has_minor_gridlines = vaxis_has_minor_gridlines
    value_axis.minor_tick_mark = tick_mark_pos_dct[vaxis_minor_tick_mark]
    value_axis.major_tick_mark = tick_mark_pos_dct[vaxis_major_tick_mark]
    value_axis.maximum_scale = vaxis_max_scale
    value_axis.minimum_scale = vaxis_min_scale
    value_axis.major_unit = vaxis_major_unit
    value_axis.minor_unit = vaxis_minor_unit
    value_axis.tick_label_position = tick_label_pos_dct[vaxis_tick_label_position]

    value_axis.visible = vaxis_visible
    if vaxis_visible:
        vaxis_tick_labels = value_axis.tick_labels
        vaxis_tick_labels.font.name = vaxis_tick_labels_font_name
        vaxis_tick_labels.font.size = Pt(vaxis_tick_labels_font_size)
        vaxis_tick_labels.font.bold = vaxis_tick_labels_font_bold
        vaxis_tick_labels.font.italic = vaxis_tick_labels_font_italic
        vaxis_tick_labels.font.color.rgb = RGBColor(*vaxis_tick_labels_font_color)
        if vaxis_tick_labels_num_format is not None:
            vaxis_tick_labels.number_format = vaxis_tick_labels_num_format
        vaxis_tick_labels.number_format_is_linked = vaxis_tick_labels_num_format_is_linked

    # set plot area properties
    plot = chart.plots[0]
    plot.vary_by_categories = plot_vary_by_cat
    plot.gap_width = plot_gap_width
    plot.overlap = plot_overlap

    plot.has_data_labels = plot_has_data_labels
    if plot_has_data_labels:
        data_labels = plot.data_labels
        data_labels.position = data_label_pos_dct[data_labels_position]
        data_labels.font.size = Pt(data_labels_font_size)
        data_labels.font.bold = data_labels_font_bold
        data_labels.font.italic = data_labels_font_italic
        data_labels.font.name = data_labels_font_name
        data_labels.font.color.rgb = RGBColor(*data_labels_font_color)
        if data_labels_num_format  is not None:
            data_labels.number_format = data_labels_num_format
        data_labels.number_format_is_linked = data_labels_num_format_is_linked

    for i, ser in enumerate(dataframe.columns):
        ser = plot.series[i]
        ser.invert_if_negative = invert_series_color_if_negative
        ser.smooth = smooth_line

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_stacked_bar_chart(
                        sld, dataframe,
                        left=284400, top=1475999, width=8582400, height=4140000,
                        chart_style=2,

                        #Legend properties
                        has_legend=True,
                        legend_position='right',
                        legend_in_layout=False,
                        legend_horz_offset=0,
                        legend_font_name="Calibri",
                        legend_font_size=10,
                        legend_font_bold=False,
                        legend_font_italic=False,
                        legend_font_color=(89,89,89),
                        legend_font_brightness=0,

                        #Category axis properties
                        caxis_visible=True,
                        caxis_tick_label_position='none',
                        caxis_tick_labels_offset =730,
                        caxis_has_major_gridlines=False,
                        caxis_has_minor_gridlines=False,
                        caxis_major_tick_mark='outside',
                        caxis_minor_tick_mark='none',
                        caxis_tick_labels_font_name = "Calibri",
                        caxis_tick_labels_font_size=10,
                        caxis_tick_labels_font_bold=False,
                        caxis_tick_labels_font_italic=False,
                        caxis_tick_labels_font_color=(89,89,89),

                        #Value axis properties
                        vaxis_visible=True,
                        vaxis_tick_label_position='low',
                        vaxis_has_major_gridlines=True,
                        vaxis_has_minor_gridlines=False,
                        vaxis_major_tick_mark='outside',
                        vaxis_minor_tick_mark='none',
                        vaxis_max_scale=1.0,
                        vaxis_min_scale=0,
                        vaxis_major_unit=0.1,
                        vaxis_minor_unit=None,
                        vaxis_tick_labels_num_format='0%',
                        vaxis_tick_labels_num_format_is_linked=False,
                        vaxis_tick_labels_font_name="Calibri",
                        vaxis_tick_labels_font_size=10,
                        vaxis_tick_labels_font_bold=True,
                        vaxis_tick_labels_font_italic=False,
                        vaxis_tick_labels_font_color=(89,89,89),

                        #Datalabel properties
                        plot_has_data_labels=True,
                        data_labels_position='center',
                        data_labels_num_format='0%',
                        data_labels_num_format_is_linked=False,
                        data_labels_font_name="Calibri",
                        data_labels_font_size=9,
                        data_labels_font_bold=False,
                        data_labels_font_italic=False,
                        data_labels_font_color=(0,0,0),

                        #Plot properties
                        invert_series_color_if_negative=False,
                        plot_gap_width=150,
                        plot_overlap=100
                        ):
    #-------------------------------------------------------------------------

    #if category labels are split from chart shape then determine the width of the textboxes and adjust chart shape accordingly.
    #textboxes in this case will take up 30% of the overall width of the chart shape. From this we can calculate the
    #width of the chart shape.

    if (caxis_visible == False) or (caxis_visible == True and str(tick_label_pos_dct[caxis_tick_label_position]) == "NONE (-4142)"):
        catwidth = percentage_of_num(30, width)
        width = width - catwidth
        left = left + catwidth

    # orientation of chart type requires that we flip/transpose the table and reverse the rows
    dataframe = dataframe.T
    dataframe = dataframe[::-1]

    # add data
    chart_data = ChartData()
    chart_data.categories = dataframe.index

    for i, col in enumerate(dataframe.columns):
        chart_data.add_series(col, (dataframe.ix[:, i].values), '0.00%')

    # add chart to slide
    x, y, cx, cy = Emu(left), Emu(top), Emu(width), Emu(height)
    graphic_frame = sld.shapes.add_chart(
        XL_CHART_TYPE.BAR_STACKED_100, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart

    # ---------------- adjust chart properties ----------------

    # chart style
    chart.chart_style = chart_style

    # set legend properties
    chart.has_legend = has_legend
    if has_legend:
        legend = chart.legend
        legend.include_in_layout = legend_in_layout
        legend.position = legend_pos_dct[legend_position]
        legend.horz_offset = legend_horz_offset
        legend.font.name = legend_font_name
        legend.font.size = Pt(legend_font_size)
        legend.font.bold = legend_font_bold
        legend.font.italic = legend_font_italic
        legend.font.color.rgb  = RGBColor(*legend_font_color)
        legend.font.color.brightness = legend_font_brightness

    # set category axis (vertical) properties
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = caxis_has_major_gridlines
    category_axis.has_minor_gridlines = caxis_has_minor_gridlines
    category_axis.major_tick_mark = tick_mark_pos_dct[caxis_major_tick_mark]
    category_axis.minor_tick_mark = tick_mark_pos_dct[caxis_minor_tick_mark]
    category_axis.tick_label_position = tick_label_pos_dct[caxis_tick_label_position]

    category_axis.visible = caxis_visible
    if caxis_visible:
        caxis_tick_labels = category_axis.tick_labels
        caxis_tick_labels.offset = caxis_tick_labels_offset
        caxis_tick_labels.font.name = caxis_tick_labels_font_name
        caxis_tick_labels.font.size = Pt(caxis_tick_labels_font_size)
        caxis_tick_labels.font.bold  = caxis_tick_labels_font_bold
        caxis_tick_labels.font.italic = caxis_tick_labels_font_italic
        caxis_tick_labels.font.color.rgb = RGBColor(*caxis_tick_labels_font_color)

    # set value axis (horizontal) properties
    value_axis = chart.value_axis
    value_axis.has_major_gridlines = vaxis_has_major_gridlines
    value_axis.has_minor_gridlines = vaxis_has_minor_gridlines
    value_axis.maximum_scale = vaxis_max_scale
    value_axis.minimum_scale = vaxis_min_scale
    value_axis.major_unit = vaxis_major_unit
    value_axis.minor_unit = vaxis_minor_unit
    value_axis.major_tick_mark = tick_mark_pos_dct[vaxis_major_tick_mark]
    value_axis.minor_tick_mark = tick_mark_pos_dct[vaxis_minor_tick_mark]
    value_axis.tick_label_position = tick_label_pos_dct[vaxis_tick_label_position]

    value_axis.visible = vaxis_visible
    if vaxis_visible:
        vaxis_tick_labels = value_axis.tick_labels
        vaxis_tick_labels.font.bold = vaxis_tick_labels_font_bold
        vaxis_tick_labels.font.size = Pt(vaxis_tick_labels_font_size)
        vaxis_tick_labels.font.italic = vaxis_tick_labels_font_italic
        vaxis_tick_labels.font.name = vaxis_tick_labels_font_name
        vaxis_tick_labels.font.color.rgb = RGBColor(*vaxis_tick_labels_font_color)
        if vaxis_tick_labels_num_format is not None:
            vaxis_tick_labels.number_format = vaxis_tick_labels_num_format
        vaxis_tick_labels.number_format_is_linked = vaxis_tick_labels_num_format_is_linked

    # set plot area properties
    plot = chart.plots[0]
    plot.has_data_labels = plot_has_data_labels
    plot.gap_width = plot_gap_width
    plot.overlap = plot_overlap

    if plot_has_data_labels:
        data_labels = plot.data_labels
        data_labels.position = data_label_pos_dct[data_labels_position]
        data_labels.font.size = Pt(data_labels_font_size)
        data_labels.font.bold = data_labels_font_bold
        data_labels.font.italic = data_labels_font_italic
        data_labels.font.name = data_labels_font_name
        data_labels.font.color.rgb = RGBColor(*data_labels_font_color)
        if data_labels_num_format is not None:
            data_labels.number_format = data_labels_num_format
        data_labels.number_format_is_linked = data_labels_num_format_is_linked

    for i, ser in enumerate(dataframe.columns):
        ser = plot.series[i]
        ser.invert_if_negative = invert_series_color_if_negative

    # generate overlay axis labels
    if (caxis_visible == False) or (caxis_visible == True and str(tick_label_pos_dct[caxis_tick_label_position]) == "NONE (-4142)"):

        cht_plot_height = get_cht_plot_height(height)
        heightPerLabel = cht_plot_height/len(dataframe.index)
        rightofchart = left + width
        txtbx_width = width / 5
        firstposition = top + get_upper_cht_plot_gap(height)

        cat_labels = dataframe.T.columns

        for i, label in enumerate(cat_labels):

            top = 0
            pointRelPos = len(cat_labels) - (i + 1)
            top = firstposition + pointRelPos * heightPerLabel

            add_textbox(sld,
                        text=label,
                        left=142875, top=top, width=rightofchart - width, height=heightPerLabel,
                        font_name=caxis_tick_labels_font_name,
                        font_size=caxis_tick_labels_font_size,
                        fit_text=False,
                        word_wrap=True,
                        font_bold=False,
                        font_color=caxis_tick_labels_font_color,
                        horizontal_alignment='right',
                        vertical_alignment='middle')

'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def chart_selector(slide, df, chart_type, *args, **kwargs):

    if chart_type == "bar":
        add_bar_chart(slide, df, *args, **kwargs)
    elif chart_type == "stacked_bar":
        add_stacked_bar_chart(slide, df, *args, **kwargs)
    elif chart_type == "column":
        add_column_chart(slide, df, *args, **kwargs)
    elif chart_type == "pie":
        add_pie_chart(slide, df, *args, **kwargs)
    elif chart_type == "line":
        add_line_chart(slide, df, *args, **kwargs)
    else:
        raise ValueError('chart type not found')


'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'
'~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~'

def add_omni_summary_slide(sld):
    #shape1
    slide_title_shp = add_textbox(sld,
                        text='Custom research',
                        left=284400, top=309600, width=8582400, height=691200,
                        font_name='Calibri',
                        font_size=36,
                        fit_text=False,
                        font_bold=False,
                        font_color=(0,0,0),
                        vertical_alignment='middle')

    #---------------------------------------------------------------------------------------------------------
    #shape2
    textbox = sld.shapes.add_textbox(Emu(343079), Emu(1477909), Emu(5411175), Emu(553998))

    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = paragraph_alignment_pos_dct['justify']

    paragraph_str = ['The charts in this presentation show your "topline" '
                     'findings but did you know that we are able to do ',
                     'so much more for you?!']

    for i, para_str in enumerate(paragraph_str):
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.color.rgb = RGBColor(89,89,89)
        if i in [0]:
            run.text = para_str
            font.bold = False
            font.size = Pt(14)
        else:
            run.text = para_str
            font.bold = True
            font.size = Pt(16)

    #---------------------------------------------------------------------------------------------------------
    #shape3
    textbox = sld.shapes.add_textbox(Emu(343079), Emu(2362766), Emu(4524484), Emu(584775))

    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = paragraph_alignment_pos_dct['justify']

    paragraph_str = ['We have ','expert specialised research teams ',
                     'who know ','your industry ', 'and ','audiences ',
                     'inside-out. ']

    for i, para_str in enumerate(paragraph_str):
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.color.rgb = RGBColor(89,89,89)
        if i in [0,2,4,6]:
            run.text = para_str
            font.bold = False
            font.size = Pt(14)
        else:
            run.text = para_str
            font.bold = True
            font.size = Pt(16)

    #---------------------------------------------------------------------------------------------------------
    #shape4
    textbox = sld.shapes.add_textbox(Emu(343078), Emu(3136878), Emu(4524485), Emu(1661993))

    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = paragraph_alignment_pos_dct['justify']

    paragraph_str = ['\nOur ','sector specialists ','(Consumer, Digital, Media and Technology, '
    'Financial Services, Public Services) combine research expertise with in-depth knowledge '
    'to help you identify and analyse your markets, as well as offer ','actionable insight ',
    'on how to best achieve your business/ organisation objectives']

    for i, para_str in enumerate(paragraph_str):
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.color.rgb = RGBColor(89,89,89)
        if i in [0,2,4]:
            run.text = para_str
            font.bold = False
            font.size = Pt(14)
        else:
            run.text = para_str
            font.bold = True
            font.size = Pt(16)

    #---------------------------------------------------------------------------------------------------------
    #shape5
    textbox = sld.shapes.add_textbox(Emu(343080), Emu(4900058), Emu(4524483), Emu(1107996))

    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.paragraphs[0]
    p.alignment = paragraph_alignment_pos_dct['justify']

    #TODO add information for footer (contact details)
    paragraph_str = ['',
                     '', '']

    for i, para_str in enumerate(paragraph_str):
        run = p.add_run()
        font = run.font
        font.name = 'Calibri'
        font.color.rgb = RGBColor(89,89,89)
        if i in [0,2]:
            run.text = para_str
            font.bold = False
            font.size = Pt(14)
        else:
            run.text = para_str
            font.bold = False
            font.size = Pt(14)
            run.hyperlink.address = para_str

    #---------------------------------------------------------------------------------------------------------
    #shape6
    img_path = 'I:\\Majeed Sahebzadha\\Images\\omnibus\\'

    img5 = sld.shapes.add_picture(img_path+'Picture5.png', Emu(6712792), Emu(4027375), Emu(2313396), Emu(1763358))
    img3 = sld.shapes.add_picture(img_path+'Picture3.png', Emu(6712792), Emu(771861), Emu(2240834), Emu(1679237))
    img4 = sld.shapes.add_picture(img_path+'Picture4.png', Emu(6070901), Emu(3350535), Emu(2291820), Emu(1702950))
    img2 = sld.shapes.add_picture(img_path+'Picture2.png', Emu(6094665), Emu(1518224), Emu(2244292), Emu(1689084))
    img1 = sld.shapes.add_picture(img_path+'Picture1.png', Emu(5238000), Emu(2494800), Emu(2307153), Emu(1740484))

    #---------------------------------------------------------------------------------------------------------
